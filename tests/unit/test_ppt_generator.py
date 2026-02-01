"""
Unit tests for PPT generator
"""

import pytest
from pathlib import Path
from pptx import Presentation

from src.generator.ppt_generator import PPTGenerator
from src.outline.models import Outline, SlideOutline, OutlineMetadata
from src.template.loader import TemplateLoader


class TestPPTGenerator:
    """Test PPTGenerator class"""

    @pytest.fixture
    def simple_outline(self):
        """Create a simple test outline"""
        slides = [
            SlideOutline(
                slide_number=1,
                layout_type="cover",
                content={"title": "Test Presentation", "subtitle": "Test Subtitle"},
            ),
            SlideOutline(
                slide_number=2,
                layout_type="toc",
                content={"title": "Contents", "items": ["Topic 1", "Topic 2"]},
            ),
            SlideOutline(
                slide_number=3,
                layout_type="content_single",
                content={"title": "Topic 1", "content": "Some content here"},
            ),
            SlideOutline(
                slide_number=4,
                layout_type="ending",
                content={"message": "Thank You!"},
            ),
        ]

        metadata = OutlineMetadata(
            total_slides=4, llm_model="test", prompt="Test presentation"
        )

        return Outline(
            outline_id="test_ppt",
            title="Test Presentation",
            template_id="business_001",
            slides=slides,
            metadata=metadata,
        )

    @pytest.fixture
    def generator(self):
        """PPT generator fixture"""
        return PPTGenerator()

    def test_generate_ppt_success(self, generator, simple_outline, tmp_path):
        """Test successful PPT generation"""
        output_path = tmp_path / "test.pptx"

        result_path = generator.generate(simple_outline, output_path)

        assert result_path.exists()
        assert result_path.suffix == ".pptx"

        # Verify PPT can be loaded
        prs = Presentation(str(result_path))
        assert len(prs.slides) == 4

    def test_generate_auto_output_path(self, generator, simple_outline, tmp_path):
        """Test PPT generation with auto-generated output path"""
        # Temporarily change output dir
        from src.utils import config

        original_output_dir = config.settings.output_dir
        config.settings.output_dir = str(tmp_path)

        try:
            result_path = generator.generate(simple_outline)
            assert result_path.exists()
            assert result_path.name == "test_ppt.pptx"
        finally:
            config.settings.output_dir = original_output_dir

    def test_generate_with_template(self, generator, simple_outline, tmp_path):
        """Test PPT generation with explicit template"""
        template_loader = TemplateLoader()
        template = template_loader.load_template("business_001")

        output_path = tmp_path / "test_with_template.pptx"
        result_path = generator.generate(simple_outline, output_path, template)

        assert result_path.exists()

    def test_generate_invalid_template(self, generator, simple_outline, tmp_path):
        """Test PPT generation with invalid template ID"""
        from src.exceptions import PPTGenerationError

        simple_outline.template_id = "nonexistent_template"
        output_path = tmp_path / "test.pptx"

        with pytest.raises(PPTGenerationError):
            generator.generate(simple_outline, output_path)

    def test_generate_different_layouts(self, generator, tmp_path):
        """Test PPT generation with different layout types"""
        slides = [
            SlideOutline(
                slide_number=1, layout_type="cover", content={"title": "Test"}
            ),
            SlideOutline(
                slide_number=2, layout_type="toc", content={"title": "TOC"}
            ),
            SlideOutline(
                slide_number=3,
                layout_type="content_two_column",
                content={
                    "title": "Two Columns",
                    "content_left": "Left",
                    "content_right": "Right",
                },
            ),
            SlideOutline(
                slide_number=4, layout_type="ending", content={"message": "End"}
            ),
        ]

        metadata = OutlineMetadata(
            total_slides=4, llm_model="test", prompt="Test"
        )

        outline = Outline(
            outline_id="test_layouts",
            title="Layout Test",
            template_id="business_001",
            slides=slides,
            metadata=metadata,
        )

        output_path = tmp_path / "layouts_test.pptx"
        result_path = generator.generate(outline, output_path)

        assert result_path.exists()

        # Verify all slides are created
        prs = Presentation(str(result_path))
        assert len(prs.slides) == 4
