"""
Style application utilities
"""

from typing import Optional
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.shapes.autoshape import Shape
from pptx.text.text import TextFrame, _Paragraph

from src.template.models import FontConfig, Placeholder
from src.utils.logger import get_logger

logger = get_logger(__name__)


class StyleApplicator:
    """Apply styles to PPT elements"""

    # Slide dimensions (16:9 aspect ratio)
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(5.625)

    @staticmethod
    def apply_font_style(paragraph: _Paragraph, font_config: FontConfig) -> None:
        """
        Apply font configuration to paragraph

        Args:
            paragraph: python-pptx paragraph object
            font_config: Font configuration
        """
        font = paragraph.font
        font.name = font_config.name
        font.size = Pt(font_config.size)
        font.bold = font_config.bold
        font.italic = font_config.italic

        if font_config.color:
            # Convert hex color to RGBColor
            r, g, b = ColorHelper.hex_to_rgb(font_config.color)
            font.color.rgb = RGBColor(r, g, b)

    @staticmethod
    def apply_text_alignment(paragraph: _Paragraph, alignment: str) -> None:
        """
        Apply text alignment to paragraph

        Args:
            paragraph: python-pptx paragraph object
            alignment: Alignment string ("left", "center", "right")
        """
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }

        if alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
        else:
            logger.warning(f"Unknown alignment: {alignment}, using left")
            paragraph.alignment = PP_ALIGN.LEFT

    @staticmethod
    def get_position(
        placeholder: Placeholder, slide_width: float, slide_height: float
    ) -> tuple:
        """
        Convert relative position to absolute inches

        Args:
            placeholder: Placeholder with relative coordinates (0-1)
            slide_width: Slide width in inches
            slide_height: Slide height in inches

        Returns:
            Tuple of (left, top, width, height) in Inches
        """
        left = Inches(placeholder.x * slide_width.inches)
        top = Inches(placeholder.y * slide_height.inches)
        width = Inches(placeholder.width * slide_width.inches)
        height = Inches(placeholder.height * slide_height.inches)

        return left, top, width, height

    @staticmethod
    def apply_background_color(slide, color_hex: str) -> None:
        """
        Apply solid background color to slide

        Args:
            slide: python-pptx slide object
            color_hex: Hex color string (e.g., "#1F4788")
        """
        try:
            background = slide.background
            fill = background.fill
            fill.solid()

            # Convert hex to RGBColor
            r, g, b = ColorHelper.hex_to_rgb(color_hex)
            fill.fore_color.rgb = RGBColor(r, g, b)

        except Exception as e:
            logger.warning(f"Failed to set background color: {e}")

    @staticmethod
    def auto_fit_text(
        text_frame: TextFrame, max_font_size: int = 44, min_font_size: int = 12
    ) -> None:
        """
        Auto-fit text by adjusting font size

        Args:
            text_frame: TextFrame to adjust
            max_font_size: Maximum font size
            min_font_size: Minimum font size
        """
        # Enable word wrap
        text_frame.word_wrap = True

        # Try to auto-fit text
        # Note: This is a basic implementation
        # More sophisticated fitting would check actual text overflow
        for paragraph in text_frame.paragraphs:
            if paragraph.font.size:
                current_size = paragraph.font.size.pt
                if current_size > max_font_size:
                    paragraph.font.size = Pt(max_font_size)
                elif current_size < min_font_size:
                    paragraph.font.size = Pt(min_font_size)


class ColorHelper:
    """Color conversion utilities"""

    @staticmethod
    def hex_to_rgb(hex_color: str) -> tuple:
        """
        Convert hex color to RGB tuple

        Args:
            hex_color: Hex color string (e.g., "#1F4788")

        Returns:
            Tuple of (r, g, b) integers
        """
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    @staticmethod
    def rgb_to_int(r: int, g: int, b: int) -> int:
        """
        Convert RGB to integer for python-pptx

        Args:
            r: Red component (0-255)
            g: Green component (0-255)
            b: Blue component (0-255)

        Returns:
            RGB integer
        """
        return (r << 16) | (g << 8) | b
