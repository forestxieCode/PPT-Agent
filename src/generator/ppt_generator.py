"""
PPT Generator - Main generator class
"""

from pathlib import Path
from typing import Optional
from pptx import Presentation

from src.outline.models import Outline
from src.template.models import Template
from src.template.loader import TemplateLoader
from src.generator.styling import StyleApplicator
from src.generator.renderers.cover import CoverRenderer
from src.generator.renderers.toc import TOCRenderer
from src.generator.renderers.content import ContentRenderer
from src.generator.renderers.ending import EndingRenderer
from src.utils.logger import get_logger
from src.utils.config import settings
from src.utils.file_utils import ensure_dir
from src.exceptions import PPTGenerationError

logger = get_logger(__name__)


class PPTGenerator:
    """PPT generator from outline and template"""

    def __init__(self, template_loader: Optional[TemplateLoader] = None):
        """
        Initialize PPT generator

        Args:
            template_loader: Template loader instance (optional)
        """
        self.template_loader = template_loader or TemplateLoader()
        self.style = StyleApplicator()

        # Initialize renderers
        self.cover_renderer = CoverRenderer(self.style)
        self.toc_renderer = TOCRenderer(self.style)
        self.content_renderer = ContentRenderer(self.style)
        self.ending_renderer = EndingRenderer(self.style)

        logger.info("Initialized PPTGenerator")

    def generate(
        self,
        outline: Outline,
        output_path: Optional[Path] = None,
        template: Optional[Template] = None,
    ) -> Path:
        """
        Generate PPT from outline

        Args:
            outline: Outline object with slides
            output_path: Output file path (optional, auto-generated if not provided)
            template: Template object (optional, loaded from outline if not provided)

        Returns:
            Path to generated PPT file

        Raises:
            PPTGenerationError: Failed to generate PPT
        """
        try:
            logger.info(f"Generating PPT: {outline.title}")

            # Load template if not provided
            if template is None:
                template = self.template_loader.load_template(outline.template_id)
                logger.debug(f"Loaded template: {template.template_name}")

            # Create presentation
            prs = Presentation()
            prs.slide_width = self.style.SLIDE_WIDTH
            prs.slide_height = self.style.SLIDE_HEIGHT

            # Generate each slide
            for slide_outline in outline.slides:
                logger.debug(
                    f"Generating slide {slide_outline.slide_number}: {slide_outline.layout_type}"
                )

                # Get layout configuration
                layout_type = slide_outline.layout_type
                if layout_type not in template.layouts:
                    raise PPTGenerationError(
                        f"Layout '{layout_type}' not found in template"
                    )

                layout = template.layouts[layout_type]

                # Add blank slide
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)

                # Render slide based on type
                self._render_slide(
                    slide, layout, slide_outline.content, template, layout_type
                )

            # Determine output path
            if output_path is None:
                output_dir = Path(settings.output_dir)
                ensure_dir(output_dir)
                filename = f"{outline.outline_id}.pptx"
                output_path = output_dir / filename

            # Ensure output directory exists
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            prs.save(str(output_path))
            logger.info(f"PPT saved to: {output_path}")

            return output_path

        except Exception as e:
            logger.error(f"Failed to generate PPT: {e}", exc_info=True)
            raise PPTGenerationError(f"PPT generation failed: {e}")

    def _render_slide(self, slide, layout, content, template, layout_type: str) -> None:
        """
        Render a single slide using appropriate renderer

        Args:
            slide: python-pptx slide object
            layout: Layout configuration
            content: Slide content
            template: Template object
            layout_type: Layout type string
        """
        # Determine which renderer to use
        if layout_type == "cover":
            self.cover_renderer.render(slide, layout, content, template)
        elif layout_type in ["toc", "table_of_contents"]:
            self.toc_renderer.render(slide, layout, content, template)
        elif layout_type == "ending":
            self.ending_renderer.render(slide, layout, content, template)
        elif layout_type.startswith("content_"):
            # All content layouts use the content renderer
            self.content_renderer.render(slide, layout, content, template)
        else:
            logger.warning(f"Unknown layout type: {layout_type}, using content renderer")
            self.content_renderer.render(slide, layout, content, template)

    def generate_from_outline_file(
        self, outline_path: Path, output_path: Optional[Path] = None
    ) -> Path:
        """
        Generate PPT from outline JSON file

        Args:
            outline_path: Path to outline JSON file
            output_path: Output file path (optional)

        Returns:
            Path to generated PPT file
        """
        from src.utils.file_utils import load_json
        from src.outline.models import Outline

        logger.info(f"Loading outline from: {outline_path}")

        # Load outline
        outline_data = load_json(outline_path)
        outline = Outline(**outline_data)

        # Generate PPT
        return self.generate(outline, output_path)
