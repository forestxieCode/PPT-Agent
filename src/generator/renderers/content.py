"""
Content slide renderer
"""

from pptx.slide import Slide
from pptx.util import Pt
from typing import Dict, Any

from src.template.models import Template, Layout
from src.generator.styling import StyleApplicator
from src.utils.logger import get_logger

logger = get_logger(__name__)


class ContentRenderer:
    """Render content slides"""

    def __init__(self, style_applicator: StyleApplicator):
        """
        Initialize content renderer

        Args:
            style_applicator: Style applicator instance
        """
        self.style = style_applicator

    def render(
        self, slide: Slide, layout: Layout, content: Dict[str, Any], template: Template
    ) -> None:
        """
        Render content slide

        Args:
            slide: python-pptx slide object
            layout: Layout configuration
            content: Content data (title, content, etc.)
            template: Template object
        """
        logger.debug(f"Rendering content slide with layout: {layout.type}")

        # Apply background
        if layout.background and layout.background.type == "solid":
            self.style.apply_background_color(slide, layout.background.color)

        # Render each placeholder
        for placeholder in layout.placeholders:
            placeholder_id = placeholder.id
            text_content = content.get(placeholder_id, "")

            if not text_content:
                # Skip optional placeholders like images
                if placeholder.type.value in ["image", "chart"]:
                    logger.debug(f"Skipping optional placeholder: {placeholder_id}")
                    continue
                logger.warning(f"No content for placeholder: {placeholder_id}")
                continue

            # Get position
            left, top, width, height = self.style.get_position(
                placeholder, self.style.SLIDE_WIDTH, self.style.SLIDE_HEIGHT
            )

            # Add text box
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Determine if this is title or body content
            is_title = placeholder_id == "title" or placeholder.type.value == "title"

            if is_title:
                # Render as simple title
                p = text_frame.paragraphs[0]
                p.text = str(text_content)

                # Apply title font
                if placeholder.font:
                    self.style.apply_font_style(p, placeholder.font)
                else:
                    theme_fonts = template.theme.fonts
                    if "heading" in theme_fonts:
                        self.style.apply_font_style(p, theme_fonts["heading"])

                if placeholder.alignment:
                    self.style.apply_text_alignment(p, placeholder.alignment)

            else:
                # Render as formatted body text
                self._render_body_content(
                    text_frame, str(text_content), placeholder, template
                )

        logger.info(f"Content slide rendered successfully")

    def _render_body_content(
        self, text_frame, content: str, placeholder, template: Template
    ) -> None:
        """
        Render body content with formatting

        Args:
            text_frame: TextFrame to render into
            content: Text content (may contain newlines and bullets)
            placeholder: Placeholder configuration
            template: Template object
        """
        # Clear default paragraph
        text_frame.clear()

        # Split content by lines
        lines = content.split("\n")

        for line in lines:
            line = line.strip()
            if not line:
                # Add empty paragraph for spacing
                text_frame.add_paragraph()
                continue

            p = text_frame.add_paragraph()

            # Check if line is a bullet point
            if line.startswith("•") or line.startswith("-") or line.startswith("*"):
                # Remove bullet marker
                text = line[1:].strip()
                p.text = text
                p.level = 0  # Bullet level

                # Set bullet format
                # Note: python-pptx doesn't support custom bullets easily
                # We'll just indent the text
            else:
                p.text = line
                p.level = 0

            # Apply spacing
            p.space_after = Pt(8)

            # Apply font style
            if placeholder.font:
                self.style.apply_font_style(p, placeholder.font)
            else:
                theme_fonts = template.theme.fonts
                if "body" in theme_fonts:
                    self.style.apply_font_style(p, theme_fonts["body"])

            # Apply alignment
            if placeholder.alignment:
                self.style.apply_text_alignment(p, placeholder.alignment)
