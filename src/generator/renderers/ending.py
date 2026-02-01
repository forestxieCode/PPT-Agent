"""
Ending slide renderer
"""

from pptx.slide import Slide
from typing import Dict, Any

from src.template.models import Template, Layout
from src.generator.styling import StyleApplicator
from src.utils.logger import get_logger

logger = get_logger(__name__)


class EndingRenderer:
    """Render ending slide"""

    def __init__(self, style_applicator: StyleApplicator):
        """
        Initialize ending renderer

        Args:
            style_applicator: Style applicator instance
        """
        self.style = style_applicator

    def render(
        self, slide: Slide, layout: Layout, content: Dict[str, Any], template: Template
    ) -> None:
        """
        Render ending slide

        Args:
            slide: python-pptx slide object
            layout: Layout configuration
            content: Content data (message, contact)
            template: Template object
        """
        logger.debug("Rendering ending slide")

        # Apply background
        if layout.background and layout.background.type == "solid":
            self.style.apply_background_color(slide, layout.background.color)

        # Render each placeholder
        for placeholder in layout.placeholders:
            placeholder_id = placeholder.id
            text_content = content.get(placeholder_id, "")

            if not text_content:
                # Use default message if not provided
                if placeholder_id == "message":
                    text_content = "感谢聆听！"
                else:
                    logger.debug(f"No content for placeholder: {placeholder_id}")
                    continue

            # Get position
            left, top, width, height = self.style.get_position(
                placeholder, self.style.SLIDE_WIDTH, self.style.SLIDE_HEIGHT
            )

            # Add text box
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Add paragraph
            p = text_frame.paragraphs[0]
            p.text = str(text_content)

            # Apply font style
            if placeholder.font:
                self.style.apply_font_style(p, placeholder.font)
            else:
                # Use title font for message, body font for contact
                theme_fonts = template.theme.fonts
                if placeholder_id == "message" and "title" in theme_fonts:
                    self.style.apply_font_style(p, theme_fonts["title"])
                elif "body" in theme_fonts:
                    self.style.apply_font_style(p, theme_fonts["body"])

            # Apply alignment
            if placeholder.alignment:
                self.style.apply_text_alignment(p, placeholder.alignment)

        logger.info("Ending slide rendered successfully")
