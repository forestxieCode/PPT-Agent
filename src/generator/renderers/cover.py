"""
Cover slide renderer
"""

from pptx.slide import Slide
from typing import Dict, Any

from src.template.models import Template, Layout
from src.generator.styling import StyleApplicator
from src.utils.logger import get_logger

logger = get_logger(__name__)


class CoverRenderer:
    """Render cover slide"""

    def __init__(self, style_applicator: StyleApplicator):
        """
        Initialize cover renderer

        Args:
            style_applicator: Style applicator instance
        """
        self.style = style_applicator

    def render(
        self, slide: Slide, layout: Layout, content: Dict[str, Any], template: Template
    ) -> None:
        """
        Render cover slide

        Args:
            slide: python-pptx slide object
            layout: Layout configuration
            content: Content data (title, subtitle, author, date)
            template: Template object
        """
        logger.debug("Rendering cover slide")

        # Apply background if configured
        if layout.background and layout.background.type == "solid":
            self.style.apply_background_color(slide, layout.background.color)

        # Render each placeholder
        for placeholder in layout.placeholders:
            placeholder_id = placeholder.id
            text_content = content.get(placeholder_id, "")

            if not text_content:
                logger.debug(f"No content for placeholder: {placeholder_id}")
                continue

            # Get position
            left, top, width, height = self.style.get_position(
                placeholder, self.style.SLIDE_WIDTH, self.style.SLIDE_HEIGHT
            )

            # Add text box
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Add paragraph
            p = text_frame.paragraphs[0]
            p.text = str(text_content)

            # Apply font style
            if placeholder.font:
                self.style.apply_font_style(p, placeholder.font)
            else:
                # Use default font from template theme
                theme_fonts = template.theme.fonts
                if placeholder.type.value == "title" and "title" in theme_fonts:
                    self.style.apply_font_style(p, theme_fonts["title"])
                elif "subtitle" in theme_fonts:
                    self.style.apply_font_style(p, theme_fonts.get("subtitle"))

            # Apply alignment
            if placeholder.alignment:
                self.style.apply_text_alignment(p, placeholder.alignment)

        logger.info("Cover slide rendered successfully")
