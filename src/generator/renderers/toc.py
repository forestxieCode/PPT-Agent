"""
Table of contents slide renderer
"""

from pptx.slide import Slide
from pptx.util import Pt
from typing import Dict, Any, List

from src.template.models import Template, Layout
from src.generator.styling import StyleApplicator
from src.utils.logger import get_logger

logger = get_logger(__name__)


class TOCRenderer:
    """Render table of contents slide"""

    def __init__(self, style_applicator: StyleApplicator):
        """
        Initialize TOC renderer

        Args:
            style_applicator: Style applicator instance
        """
        self.style = style_applicator

    def render(
        self, slide: Slide, layout: Layout, content: Dict[str, Any], template: Template
    ) -> None:
        """
        Render table of contents slide

        Args:
            slide: python-pptx slide object
            layout: Layout configuration
            content: Content data (title, items)
            template: Template object
        """
        logger.debug("Rendering TOC slide")

        # Apply background
        if layout.background and layout.background.type == "solid":
            self.style.apply_background_color(slide, layout.background.color)

        # Render each placeholder
        for placeholder in layout.placeholders:
            placeholder_id = placeholder.id

            # Get position
            left, top, width, height = self.style.get_position(
                placeholder, self.style.SLIDE_WIDTH, self.style.SLIDE_HEIGHT
            )

            # Add text box
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            if placeholder_id == "title":
                # Render title
                title_text = content.get("title", "目录")
                p = text_frame.paragraphs[0]
                p.text = title_text

                # Apply font style
                if placeholder.font:
                    self.style.apply_font_style(p, placeholder.font)
                else:
                    theme_fonts = template.theme.fonts
                    if "title" in theme_fonts:
                        self.style.apply_font_style(p, theme_fonts["title"])

                if placeholder.alignment:
                    self.style.apply_text_alignment(p, placeholder.alignment)

            elif placeholder_id == "items":
                # Render list items
                items = content.get("items", [])
                if not items:
                    logger.warning("No TOC items provided")
                    continue

                self._render_toc_items(text_frame, items, placeholder, template)

        logger.info("TOC slide rendered successfully")

    def _render_toc_items(
        self, text_frame, items: List[str], placeholder, template: Template
    ) -> None:
        """
        Render TOC items as numbered list

        Args:
            text_frame: TextFrame to render into
            items: List of TOC items
            placeholder: Placeholder configuration
            template: Template object
        """
        # Clear default paragraph
        text_frame.clear()

        for i, item in enumerate(items, 1):
            p = text_frame.add_paragraph()
            p.text = f"{i}. {item}"
            p.level = 0
            p.space_after = Pt(12)  # Space after each item

            # Apply font style
            if placeholder.font:
                self.style.apply_font_style(p, placeholder.font)
            else:
                theme_fonts = template.theme.fonts
                if "body" in theme_fonts:
                    self.style.apply_font_style(p, theme_fonts["body"])

            # Apply alignment
            if placeholder.alignment:
                self.style.apply_text_alignment(p, placeholder.alignment)
